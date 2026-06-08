import io
import json
import logging
import re

import fitz  # pymupdf
from google.genai import types

from .clients import claude, gemini, STYLE_RULES, supabase, track_claude, track_gemini, track_gemini_embed
from .retry import QuotaExhausted, transient

log = logging.getLogger(__name__)

READ_PROMPT = (
    "Transcribe this page exactly. Keep the wording. "
    "For any math, write it in LaTeX. "
    "For any diagram or figure, give a short description inside [square brackets]."
)

OUTLINE_PROMPT = (
    "Read this study material and produce an outline of it. "
    "List the topics and sub-topics in the order they appear, "
    "like a table of contents. Keep it plain and short.\n\n"
    "OMIT non-content sections that exist for bookkeeping rather than "
    "teaching. Drop: front matter, foreword, preface, dedication, "
    "copyright page, table of contents, list of figures/tables, "
    "acknowledgments, about-the-author, glossary, index, "
    "bibliography, references, works cited, errata, colophon, and any "
    "appendix that is purely a lookup table (units, abbreviations, "
    "contact information). If a section is genuinely instructional, "
    "keep it even if it sits in the front matter.\n\n"
)

MODEL_FAST = "gemini-2.5-flash-lite"
MODEL_STRONG = "gemini-2.5-flash"

# Image extensions handled by the Gemini OCR path. Anything else routes to
# the format-specific text extractors (.docx/.pptx/.txt/.md) or PyMuPDF (.pdf).
IMAGE_EXTS = (".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".heic", ".heif")

# Map an image file extension to the mime type Gemini Vision expects. Kept for
# reference / callers that already know the extension. The OCR path no longer
# threads a filename through read_image (see _sniff_mime), so a JPEG/WEBP/HEIC
# upload isn't mislabelled as image/png even when the filename is absent.
EXT_MIME = {
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".png": "image/png",
    ".webp": "image/webp",
    ".heic": "image/heic",
    ".heif": "image/heic",
    ".gif": "image/gif",
    ".bmp": "image/bmp",
}


def _sniff_mime(img_bytes: bytes) -> str:
    """Mime type for an image, derived from its leading magic bytes. Used by
    read_image so a JPEG/WEBP/HEIC upload isn't mislabelled as image/png
    (Gemini is lenient but the correct mime avoids subtle decode quirks).
    PDF page scans and PyMuPDF crops are always rendered as PNG, so the
    default is image/png. Sniffing the bytes rather than threading a filename
    through keeps read_image's call signature to (img_bytes, model, ctx)."""
    b = img_bytes or b""
    if b.startswith(b"\xff\xd8\xff"):
        return "image/jpeg"
    if b.startswith(b"\x89PNG\r\n\x1a\n") or b.startswith(b"\x89PNG"):
        return "image/png"
    if b.startswith(b"GIF87a") or b.startswith(b"GIF89a"):
        return "image/gif"
    if len(b) >= 12 and b[:4] == b"RIFF" and b[8:12] == b"WEBP":
        return "image/webp"
    if len(b) >= 12 and b[4:8] == b"ftyp" and b[8:12] in (
        b"heic", b"heix", b"hevc", b"heim", b"heis", b"hevx", b"mif1", b"msf1",
    ):
        return "image/heic"
    if b.startswith(b"BM"):
        return "image/bmp"
    return "image/png"

# A document we couldn't get any text out of is a hard failure with an
# actionable message rather than a silent empty doc.
EMPTY_DOC_MESSAGE = (
    "We couldn't read any text from this file. If it is a very large PDF, "
    "try uploading a single chapter or a smaller file."
)

# Shown when Gemini's daily/per-minute quota runs out mid-ingest. Tells the
# user the concrete next step rather than a bare "quota reached".
QUOTA_MESSAGE = (
    "Gemini API quota reached. Try again in a few minutes or after midnight, "
    "or upload a single chapter / smaller file to stay under the limit."
)

# Shown when the user's chapter request couldn't be parsed into a number.
# `{label}` is the raw string the user typed.
CHAPTER_UNPARSEABLE_MESSAGE = (
    'Could not read the chapter "{label}". Use a number like 5 or a roman '
    "numeral like V, or choose Whole book."
)

# Shown when the chapter parsed fine but no matching chapter could be located
# in the PDF (no TOC entry and no text-layer heading). `{n}` is the number.
CHAPTER_NOT_FOUND_MESSAGE = (
    "Could not find Chapter {n} in this file. It may be a scan with no text "
    "layer or have no clear chapters. Try Whole book."
)


DIAGRAM_DETECT_PROMPT = (
    "Look at this study-page image. Find every distinct diagram, "
    "illustration, photograph, chart, anatomical figure, structural "
    "formula, schematic, map, or other visual figure that is NOT just "
    "text.\n\n"
    "IGNORE: blocks of text, headings, page numbers, footers, headers, "
    "decorative borders, tables of pure text data, the page background.\n\n"
    "For each visual figure, return its bounding box as normalized "
    "coordinates where x and y are the TOP-LEFT corner and w / h are "
    "the width / height, all in fractions of the page (0.0 to 1.0).\n\n"
    "Return ONLY a JSON object: {\"figures\":[{\"x\":0.10,\"y\":0.42,"
    "\"w\":0.35,\"h\":0.30}, ...]}. Empty array if the page is text-only."
)


# A bracketed span counts as "figure-shaped" (worth a vision call) only if
# its inner text reads like a real figure description: several words, or a
# figure keyword. A bare "[1]" citation or "[a]" label does not.
_BRACKET_RE = re.compile(r"\[([^\[\]]+)\]")
_FIGURE_KEYWORD_RE = re.compile(
    r"\b(figure|diagram|illustration|photo|chart|graph|structure|"
    r"schematic|map|fig)\b",
    re.IGNORECASE,
)


def _has_phrase_shaped_bracket(text: str) -> bool:
    """True if `text` contains a bracketed span whose inner text has >= 3
    words OR contains a figure keyword. Used to gate the (expensive) vision
    diagram-detection call on scanned pages: OCR encodes figures as
    [bracketed descriptions], so a page with no figure-shaped bracket has
    no diagram worth cropping. Biased permissive — a single keyword is
    enough."""
    for m in _BRACKET_RE.finditer(text):
        inner = m.group(1).strip()
        if not inner:
            continue
        if _FIGURE_KEYWORD_RE.search(inner):
            return True
        if len(inner.split()) >= 3:
            return True
    return False


@transient()
def _detect_page_diagrams(page, ctx=None) -> list[bytes]:
    """Vision-detect real diagrams on a scanned PDF page and return
    cropped PNG bytes for each. Used in place of `page.get_images()` on
    OCR'd pages, where PyMuPDF's only "image" is the whole-page scan
    (which would leak the question's source text if used as a figure).

    Pipeline:
      1. Render the full page at 150 dpi.
      2. Ask Gemini Vision for the bounding boxes of every non-text
         visual region on the page.
      3. For each bounding box, ask PyMuPDF to render JUST that clip of
         the original page (still at 150 dpi) → that becomes the figure.

    Returns an empty list on any failure so the caller can fall back to
    treating the page as text-only.
    """
    page_png = page.get_pixmap(dpi=150).tobytes("png")
    try:
        resp = track_gemini(
            "detect_page_diagrams",
            model=MODEL_FAST,
            contents=[
                types.Part.from_bytes(data=page_png, mime_type="image/png"),
                DIAGRAM_DETECT_PROMPT,
            ],
            ctx=ctx,
        )
        raw = (resp.text or "").strip()
    except Exception as e:
        log.info("diagram detect vision call failed: %s: %s",
                 type(e).__name__, e)
        return []

    # Strip markdown fences Gemini sometimes wraps JSON in.
    if raw.startswith("```"):
        raw = raw.split("```")[1]
        if raw.startswith("json"):
            raw = raw[4:]
    raw = raw.strip()
    if not raw:
        return []
    try:
        # Tolerate both `[...]` and `{"figures":[...]}` shapes.
        parsed = json.loads(raw)
    except json.JSONDecodeError:
        return []
    if isinstance(parsed, dict):
        bboxes = parsed.get("figures") or parsed.get("boxes") or []
    elif isinstance(parsed, list):
        bboxes = parsed
    else:
        bboxes = []
    if not isinstance(bboxes, list):
        return []

    out: list[bytes] = []
    rect = page.rect
    for box in bboxes:
        if not isinstance(box, dict):
            continue
        try:
            x = float(box.get("x", 0))
            y = float(box.get("y", 0))
            w = float(box.get("w", 0))
            h = float(box.get("h", 0))
        except (TypeError, ValueError):
            continue
        # Sanity-check the box: in-bounds and not too tiny.
        if not (0 <= x < 1 and 0 <= y < 1 and 0 < w <= 1 and 0 < h <= 1):
            continue
        if w * h < 0.02:  # less than 2% of page area — likely an icon
            continue
        clip = fitz.Rect(
            rect.x0 + x * rect.width,
            rect.y0 + y * rect.height,
            rect.x0 + (x + w) * rect.width,
            rect.y0 + (y + h) * rect.height,
        )
        try:
            pix = page.get_pixmap(dpi=150, clip=clip)
            if pix.width < 80 or pix.height < 80:
                continue
            out.append(pix.tobytes("png"))
        except Exception:
            continue
    return out


@transient()
def read_image(img_bytes: bytes, model: str = MODEL_FAST, ctx=None) -> str:
    """OCR an image with Gemini Vision. The mime type is sniffed from the
    image's leading magic bytes (jpg/webp/heic/...), defaulting to image/png
    for PDF page renders and crops."""
    resp = track_gemini(
        "ocr_page",
        model=model,
        contents=[
            types.Part.from_bytes(data=img_bytes, mime_type=_sniff_mime(img_bytes)),
            READ_PROMPT,
        ],
        ctx=ctx,
    )
    # Gemini can return None text for a blank/illegible page; never propagate
    # None into chunking (read_image is typed -> str).
    return resp.text or ""


def read_image_strong(img_bytes: bytes) -> str:
    """For handwriting and heavy math — uses gemini-2.5-flash, not flash-lite."""
    return read_image(img_bytes, model=MODEL_STRONG)


PAGE_TEXT_THRESHOLD = 200  # chars per page that count as "real text"


# =========================================================================
# Chapter selection (CHEAP — no LLM/OCR). Lets the user ingest only one
# chapter of a textbook PDF to cut API cost. Chapter detection uses the
# PDF's table-of-contents bookmarks first, then a conservative scan of the
# free text layer, never a vision/LLM call.
# =========================================================================

# Number words 1..30, including the hyphenated forms in 21..29. Lookup is
# case-insensitive (the caller lowercases) and tolerates a space instead of
# the hyphen ("twenty one") as well as the hyphenated form ("twenty-one").
_NUMBER_WORDS: dict[str, int] = {
    "one": 1, "two": 2, "three": 3, "four": 4, "five": 5,
    "six": 6, "seven": 7, "eight": 8, "nine": 9, "ten": 10,
    "eleven": 11, "twelve": 12, "thirteen": 13, "fourteen": 14,
    "fifteen": 15, "sixteen": 16, "seventeen": 17, "eighteen": 18,
    "nineteen": 19, "twenty": 20, "thirty": 30,
}
for _tens_word, _tens_val in (("twenty", 20),):
    for _ones_word, _ones_val in (
        ("one", 1), ("two", 2), ("three", 3), ("four", 4), ("five", 5),
        ("six", 6), ("seven", 7), ("eight", 8), ("nine", 9),
    ):
        _NUMBER_WORDS[f"{_tens_word}-{_ones_word}"] = _tens_val + _ones_val

# Leading keywords that introduce a chapter label ("chapter 5", "ch. V",
# "unit 3", "section II"). Stripped before parsing the number that follows.
_CHAPTER_KEYWORDS = ("chapter", "chap", "ch", "unit", "section", "part", "lesson")
_CHAPTER_KEYWORD_RE = re.compile(
    r"^(?:" + "|".join(_CHAPTER_KEYWORDS) + r")\b\.?\s*",
    re.IGNORECASE,
)

_ROMAN_RE = re.compile(r"^[ivxlcdm]+$", re.IGNORECASE)
_ROMAN_VALUES = {"i": 1, "v": 5, "x": 10, "l": 50, "c": 100, "d": 500, "m": 1000}

# Real textbook chapters numbered with roman numerals stay small. Capping the
# roman interpretation here keeps an ordinary all-roman-letter English word
# ("Mix" -> MIX -> 1009, "Did", "Civic", "Mild") from being misread as a
# chapter number. Digit and number-word parsing are unaffected by this cap.
_MAX_ROMAN_CHAPTER = 99


def _roman_to_int(s: str) -> int | None:
    """Parse a roman numeral (case-insensitive) to an int, or None if it is
    not a well-formed roman numeral. Validates by round-tripping: a malformed
    string like "iiii" or "vx" canonicalizes to a different numeral, so we
    reject it. Used only for small chapter numbers, but the full algorithm is
    cheap and keeps "XIV" etc. correct."""
    s = (s or "").strip().lower()
    if not s or not _ROMAN_RE.match(s):
        return None
    total = 0
    prev = 0
    for ch in reversed(s):
        val = _ROMAN_VALUES[ch]
        if val < prev:
            total -= val
        else:
            total += val
            prev = val
    if total <= 0:
        return None
    # Reject malformed numerals (e.g. "iiii", "vx") by canonical round-trip.
    if _int_to_roman(total) != s:
        return None
    return total


def _int_to_roman(n: int) -> str:
    """Lowercase canonical roman numeral for n (>0). Used to validate that a
    parsed roman numeral was well-formed."""
    table = (
        (1000, "m"), (900, "cm"), (500, "d"), (400, "cd"),
        (100, "c"), (90, "xc"), (50, "l"), (40, "xl"),
        (10, "x"), (9, "ix"), (5, "v"), (4, "iv"), (1, "i"),
    )
    out = []
    for val, sym in table:
        while n >= val:
            out.append(sym)
            n -= val
    return "".join(out)


def _word_to_int(s: str) -> int | None:
    """Parse an English number word (1..30, incl. hyphenated 21..29) to an
    int, or None. Tolerates a space instead of the hyphen ("twenty one")."""
    s = (s or "").strip().lower()
    if not s:
        return None
    if s in _NUMBER_WORDS:
        return _NUMBER_WORDS[s]
    # "twenty one" -> "twenty-one"
    return _NUMBER_WORDS.get(s.replace(" ", "-"))


def parse_chapter_request(s: str) -> int | None:
    """Parse a user- or TOC-supplied chapter label to a chapter number.

    Accepts plain digits ("5", "12"); roman numerals in either case
    ("V", "iv", "Xii"); number words 1..30 including hyphenated
    ("five", "twenty-one"); and any of those after a leading keyword
    ("chapter 5", "ch. V", "chap five", "unit 3", "section II"). Surrounding
    punctuation and case are ignored. Returns None when no chapter number is
    parseable.

    Pure and CHEAP — no API. Used both for the user's "chapter" field and to
    read a number out of a TOC bookmark title."""
    if not s:
        return None
    label = s.strip()
    if not label:
        return None
    # Strip a leading keyword ("chapter", "ch.", "unit", ...) if present.
    label = _CHAPTER_KEYWORD_RE.sub("", label, count=1).strip()
    # Trim surrounding punctuation/whitespace (e.g. "5." or "(V)" or "- five").
    label = label.strip(" \t\r\n.,:;-–—()[]{}\"'")
    if not label:
        return None
    # Plain digits: take the leading integer (handles "5", "12", and a TOC
    # title like "5 Cell Structure" where text trails the number).
    m = re.match(r"^(\d+)\b", label)
    if m:
        try:
            n = int(m.group(1))
        except ValueError:
            return None
        # Chapters are 1-based; "0"/"00" is not a valid chapter.
        return n if n >= 1 else None
    tokens = label.split()
    # Leading token, with its own trailing/leading punctuation stripped so a
    # TOC title like "XIV: Cells" or "V. Introduction" still parses.
    first = tokens[0].strip(".,:;-–—()[]{}\"'") if tokens else label
    # Roman numeral (a standalone token like "V" or "XIV"). Capped at a sane
    # chapter ceiling so an ordinary all-roman-letter word ("Mix", "Did",
    # "Civic") isn't misread as a huge roman numeral.
    roman = _roman_to_int(first)
    if roman is not None and roman <= _MAX_ROMAN_CHAPTER:
        return roman
    # Number word. Try a two-token hyphen-equivalent first ("twenty one"),
    # then the leading single token ("five", "twenty-one").
    if len(tokens) >= 2:
        two = _word_to_int(f"{tokens[0]}-{tokens[1]}")
        if two is not None:
            return two
    word = _word_to_int(first)
    if word is not None:
        return word
    return None


# A TOC title or text-layer line "looks like a chapter heading" when it
# starts with a digit, a roman numeral token, or the word chapter/unit/etc.
# This gates which TOC entries are considered chapter boundaries (so a
# "Preface" or "Index" bookmark never counts) and which text lines may be a
# heading in the fallback scan.
_CHAPTER_HEADING_RE = re.compile(
    r"^\s*(?:" + "|".join(_CHAPTER_KEYWORDS) + r")\b",
    re.IGNORECASE,
)
_LEADING_NUMBER_RE = re.compile(r"^\s*\d")
_LEADING_ROMAN_RE = re.compile(r"^\s*[ivxlcdm]+\b", re.IGNORECASE)


def _looks_chapter_like(title: str) -> bool:
    """Whether a TOC bookmark title reads like a chapter (vs. front-matter
    such as 'Preface' or back-matter like 'Index'). True when it starts with
    the word chapter/unit/etc., a digit, or a roman-numeral token. The
    roman test also requires parse_chapter_request to read a number, so a
    word like 'Mix' (all-roman letters) doesn't false-positive."""
    if not title:
        return False
    if _CHAPTER_HEADING_RE.match(title):
        return True
    if _LEADING_NUMBER_RE.match(title):
        return True
    if _LEADING_ROMAN_RE.match(title):
        first = title.strip().split()[0].strip(".,:;-–—()[]{}")
        return _roman_to_int(first) is not None
    return False


def _chapter_range_from_toc(doc, n: int) -> tuple[int, int] | None:
    """Locate chapter `n` from the PDF's table-of-contents bookmarks.

    doc.get_toc() yields [level, title, page] entries (page is 1-indexed).
    We collect the chapter-like entries (prefer level-1 bookmarks; fall back
    to any level when no level-1 entry looks chapter-like), parse a number
    out of each title, and find the one whose number == n. Its page is the
    start; the end is the page before the next chapter-like entry, or the
    last page of the document. Returns None when the TOC has no usable
    chapter-`n` entry. CHEAP — no API."""
    try:
        toc = doc.get_toc(simple=True) or []
    except Exception as e:
        log.warning("get_toc failed during chapter detect: %s", e)
        return None
    if not toc:
        return None
    total = len(doc)

    def _collect(level_filter) -> list[tuple[int, int]]:
        # Returns [(chapter_number, start_page), ...] in TOC order for the
        # chapter-like entries matching the level filter.
        out: list[tuple[int, int]] = []
        for entry in toc:
            try:
                level, title, page = entry[0], entry[1], entry[2]
            except (IndexError, TypeError):
                continue
            if level_filter is not None and level != level_filter:
                continue
            if not _looks_chapter_like(title):
                continue
            num = parse_chapter_request(title)
            if num is None:
                continue
            try:
                start = int(page)
            except (TypeError, ValueError):
                continue
            if start < 1:
                continue
            out.append((num, min(start, total)))
        return out

    entries = _collect(1)
    if not any(num == n for num, _ in entries):
        # No level-1 match — widen to chapter-like entries at any level.
        entries = _collect(None)
    if not entries:
        return None

    for i, (num, start) in enumerate(entries):
        if num != n:
            continue
        # End = page before the next chapter-like entry that starts strictly
        # after this one, else the last page.
        end = total
        for next_num, next_start in entries[i + 1:]:
            if next_start > start:
                end = next_start - 1
                break
        if end < start:
            end = start
        return (start, min(end, total))
    return None


# A standalone "chapter N" heading in the text layer: optionally a keyword,
# then the number as digits / roman / word, with little else on the line.
_TEXT_HEADING_MAX_LEN = 60  # a heading is a SHORT line, not buried prose


def _line_chapter_number(line: str) -> int | None:
    """If a text-layer line reads like a standalone chapter heading, return
    its chapter number, else None. Conservative: the line must be short and
    must begin with a chapter keyword or a bare number/roman/word token, so
    a number buried mid-paragraph never counts."""
    stripped = (line or "").strip()
    if not stripped or len(stripped) > _TEXT_HEADING_MAX_LEN:
        return None
    has_keyword = bool(_CHAPTER_HEADING_RE.match(stripped))
    # Without a keyword, require the line to START with a number/roman/word
    # token AND be a short heading (avoids matching ordinary sentences).
    if not has_keyword and not (
        _LEADING_NUMBER_RE.match(stripped) or _LEADING_ROMAN_RE.match(stripped)
    ):
        return None
    return parse_chapter_request(stripped)


def _chapter_range_from_text(doc, n: int) -> tuple[int, int] | None:
    """Fallback chapter locator: scan each page's free text layer for a
    standalone 'chapter n' heading. start = first page with an `n` heading;
    end = the page before the first `n+1` heading after start, else the last
    page. Conservative (short heading lines only) to avoid false positives.
    Returns None when no `n` heading is found. CHEAP — text layer only, no
    OCR / API."""
    total = len(doc)
    start: int | None = None
    end = total
    for i in range(total):
        try:
            text = doc[i].get_text()
        except Exception:
            continue
        if not text:
            continue
        page_no = i + 1
        for line in text.splitlines():
            num = _line_chapter_number(line)
            if num is None:
                continue
            if start is None:
                if num == n:
                    start = page_no
                    break
            else:
                # Already inside chapter n; the next chapter heading (n+1 or
                # any higher number) ends it.
                if num >= n + 1 and page_no > start:
                    return (start, max(page_no - 1, start))
        # A later line on the start page can't shorten the range; keep scanning
        # subsequent pages for the next chapter heading.
    if start is None:
        return None
    return (start, min(end, total))


def find_chapter_page_range(doc, n: int) -> tuple[int, int] | None:
    """1-indexed inclusive (start, end) page range of chapter `n` in an open
    PyMuPDF document, or None when it can't be located. CHEAP — no API/OCR.

    PRIMARY: the PDF's table-of-contents bookmarks (doc.get_toc()).
    FALLBACK: a conservative scan of the free text layer for a standalone
    'chapter n' heading. Returns None if neither locates the chapter."""
    found = _chapter_range_from_toc(doc, n)
    if found is not None:
        return found
    return _chapter_range_from_text(doc, n)


# =========================================================================
# Format-specific text extractors. Each returns a list of (page_number,
# text). Office-doc figures are out of scope — text only. PyMuPDF (.pdf)
# and the image OCR path keep their own routines because they also feed
# the figure-extraction stage below.
# =========================================================================

def _ext(filename: str) -> str:
    """Lowercased file extension including the dot, or '' if none."""
    dot = filename.rfind(".")
    return filename[dot:].lower() if dot != -1 else ""


def extract_docx(file_bytes: bytes) -> list[tuple[int, str]]:
    """python-docx has no page concept, so the whole document is one
    'page'. Joins non-empty paragraphs with newlines."""
    import docx  # python-docx
    doc = docx.Document(io.BytesIO(file_bytes))
    paras = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    return [(1, "\n".join(paras))]


def extract_pptx(file_bytes: bytes) -> list[tuple[int, str]]:
    """One slide -> one page. Pulls text from every shape that has a
    text frame, in slide order."""
    from pptx import Presentation  # python-pptx
    prs = Presentation(io.BytesIO(file_bytes))
    pages: list[tuple[int, str]] = []
    for i, slide in enumerate(prs.slides):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        parts.append(line)
        pages.append((i + 1, "\n".join(parts)))
    return pages


def extract_text_file(file_bytes: bytes) -> list[tuple[int, str]]:
    """.txt / .md — plain utf-8 decode, single page."""
    return [(1, file_bytes.decode("utf-8", errors="replace"))]


def extract_pages(file_bytes: bytes, filename: str,
                  doc_id: str | None = None,
                  page_range: tuple[int, int] | None = None
                  ) -> tuple[list[tuple[int, str]], str, set[int]]:
    """Return (pages, source_type, ocr_page_numbers).

    `pages` is a list of (page_number, text). For PDFs, each page is
    decided independently: if it has substantial selectable text (≥
    PAGE_TEXT_THRESHOLD chars), we use it directly; otherwise we OCR it.
    This handles mixed PDFs (mostly text with a few image pages)
    correctly and catches image-heavy PDFs with a thin text overlay.

    `ocr_page_numbers` is the set of 1-indexed page numbers that had to
    be OCR'd. The caller uses this to skip figure extraction on those
    pages — an OCR'd page has no real text layer, so the only thing
    PyMuPDF can extract from it is the page-scan image itself, which
    would leak the answer if shown as a "figure" during a test.

    `page_range` (PDF only) restricts processing to the 1-indexed inclusive
    page span (start, end) — used by the chapter-extract path so only one
    chapter is OCR'd / embedded. The REAL pdf page numbers are kept (pages
    are not renumbered) so figures, citations and resume cursors stay
    correct. Non-PDF formats ignore `page_range`.

    When `doc_id` is provided, emits `documents.progress` updates per
    page so the frontend's ingest screen can show OCR progress on long
    scanned PDFs (otherwise this stage looks frozen for minutes).
    """
    ext = _ext(filename)
    if ext == ".docx":
        return extract_docx(file_bytes), "docx", set()
    if ext == ".pptx":
        return extract_pptx(file_bytes), "pptx", set()
    if ext in (".txt", ".md"):
        return extract_text_file(file_bytes), "text", set()
    if ext != ".pdf":
        # Single-image upload: no concept of "page scan vs embedded
        # diagram" applies — return an empty exclude set. Mime comes from
        # the upload's extension so a JPEG/WEBP/HEIC isn't read as PNG.
        return [(1, read_image(file_bytes, ctx={"doc_id": doc_id}))], "image", set()

    doc = fitz.open(stream=file_bytes, filetype="pdf")
    total = len(doc)
    # Clamp the requested range into the document's real page bounds; None
    # means the whole document (unchanged behavior).
    if page_range is not None:
        lo = max(1, page_range[0])
        hi = min(total, page_range[1])
    else:
        lo, hi = 1, total
    pages = []
    ocr_page_numbers: set[int] = set()
    for i, page in enumerate(doc):
        page_no = i + 1
        if page_no < lo or page_no > hi:
            continue
        if doc_id is not None and (page_no == lo or page_no % 5 == 0 or page_no == hi):
            _set_progress(doc_id, f"reading page {page_no} of {total}")
        direct = page.get_text()
        if len(direct.strip()) >= PAGE_TEXT_THRESHOLD:
            pages.append((page_no, direct))
        else:
            img = page.get_pixmap(dpi=150).tobytes("png")
            pages.append((page_no, read_image(img, ctx={"doc_id": doc_id})))
            ocr_page_numbers.add(page_no)

    # source_type is decided over the pages we actually processed (the
    # restricted span for a chapter), matching the whole-doc behavior.
    processed = len(pages)
    source_type = "scanned" if processed and len(ocr_page_numbers) > processed / 2 else "pdf_text"
    return pages, source_type, ocr_page_numbers


def extract_image_pages(files: list[tuple[bytes, str]],
                        doc_id: str | None = None) -> list[tuple[int, str]]:
    """OCR a set of image uploads into ordered pages for ONE document:
    [(1, ocr(img1)), (2, ocr(img2)), ...]. Order follows the upload
    order. Emits per-page progress when `doc_id` is given."""
    total = len(files)
    pages: list[tuple[int, str]] = []
    for i, (img_bytes, name) in enumerate(files):
        if doc_id is not None:
            _set_progress(doc_id, f"reading page {i + 1} of {total}")
        pages.append((i + 1, read_image(img_bytes, ctx={"doc_id": doc_id})))
    return pages


def extract_page_images(file_bytes: bytes, filename: str,
                        skip_pages: set[int] | None = None,
                        only_page: int | None = None,
                        doc_id: str | None = None
                        ) -> dict[int, list[bytes]]:
    """Return {1-indexed-page: [png_bytes, ...]} of embedded images in a PDF.

    Used so the lesson screen can render the diagram alongside its
    `[bracketed]` description. The bracket descriptions stay in chunks as
    text (so Claude can still reason over them); the rendered image just
    sits next to the chunk it belongs to.

    `skip_pages` is the set of page numbers (1-indexed) that were OCR'd
    (no real text layer). On those pages PyMuPDF's `page.get_images()`
    would only find the whole-page scan, which leaks question text. We
    handle those pages differently: send the page to Gemini Vision and
    ask for the bounding boxes of every distinct diagram / illustration
    / photograph, then crop just those regions from the original page
    rendering. That way scanned textbooks with real diagrams (anatomy,
    chemistry, etc) still get the diagrams as figures, but the
    surrounding text never makes it into a `figure_path`.

    `only_page` restricts extraction to a single 1-indexed page (used by
    the page-by-page ingest loop so resumed ingests don't re-scan the
    whole PDF). When None, every page is processed. The vision-call gate
    on no-figure pages lives in `_page_figures`, which calls this.

    For text-extracted pages (NOT in skip_pages), uses the normal
    PyMuPDF embedded-image path.

    Skips images smaller than 80 x 80 px to filter out icons / bullets.
    Returns an empty dict for non-PDFs (single-image uploads already are
    the figure).
    """
    if not filename.lower().endswith(".pdf"):
        return {}
    skip = skip_pages or set()
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    out: dict[int, list[bytes]] = {}
    for i, page in enumerate(doc):
        page_no = i + 1
        if only_page is not None and page_no != only_page:
            continue
        if page_no in skip:
            # Scanned page — use vision-detected diagram regions.
            try:
                cropped = _detect_page_diagrams(page, ctx={"doc_id": doc_id})
            except Exception as e:
                log.warning(
                    "diagram detect failed for page=%s: %s: %s",
                    page_no, type(e).__name__, e)
                cropped = []
            if cropped:
                out.setdefault(page_no, []).extend(cropped)
            continue
        # Text-extracted page — pull embedded image objects.
        for img in page.get_images(full=True):
            xref = img[0]
            try:
                pix = fitz.Pixmap(doc, xref)
                if pix.width < 80 or pix.height < 80:
                    pix = None
                    continue
                if pix.n - pix.alpha > 3:  # CMYK -> convert to RGB
                    pix = fitz.Pixmap(fitz.csRGB, pix)
                png = pix.tobytes("png")
                out.setdefault(page_no, []).append(png)
                pix = None
            except Exception as e:
                log.warning("figure extract failed for page=%s xref=%s: %s",
                            page_no, xref, e)
    return out


def _page_figures(file_bytes: bytes, filename: str, page_no: int,
                  is_ocr_page: bool, ocr_text: str,
                  doc_id: str | None = None) -> list[bytes]:
    """Extract the figure PNGs for a SINGLE page. On scanned (OCR'd)
    pages this runs the vision diagram detector, but only when the page's
    OCR text carries a figure-shaped bracket — otherwise the expensive
    vision call is skipped (and logged). On text-layer PDF pages it pulls
    embedded image objects. Returns [] for non-PDFs and on any failure."""
    if not filename.lower().endswith(".pdf"):
        return []
    if is_ocr_page and not _has_phrase_shaped_bracket(ocr_text or ""):
        log.info("skipping diagram detect on page=%s for doc_id=%s "
                 "(no figure-shaped bracket)", page_no, doc_id)
        return []
    try:
        page_pngs = extract_page_images(
            file_bytes, filename,
            skip_pages={page_no} if is_ocr_page else set(),
            only_page=page_no,
            doc_id=doc_id,
        )
    except Exception:
        log.exception("figure extract failed for page=%s doc_id=%s",
                      page_no, doc_id)
        return []
    return page_pngs.get(page_no, [])


_CONTROL_CHARS_RE = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]")


def _sanitize_text(s: str | None) -> str:
    """Strip NULL bytes and other unprintable controls. Postgres TEXT columns
    reject \\u0000, and PyMuPDF occasionally leaks them through when a page
    has embedded forms or scanner artifacts. Keep \\t \\n \\r (printable
    whitespace) so layout-ish formatting survives.

    Tolerates None/empty (e.g. a blank page where Gemini OCR returns no text):
    returns "" so the page becomes a zero-chunk page instead of crashing ingest."""
    if not s:
        return ""
    return _CONTROL_CHARS_RE.sub("", s)


def chunk_pages(pages: list[tuple[int, str]], size: int = 800, overlap: int = 100
                ) -> list[tuple[str, int]]:
    """Chunk per page so page_number is reliable. Returns (text, page_number)."""
    chunks = []
    for page_num, text in pages:
        words = _sanitize_text(text).split()
        start = 0
        while start < len(words):
            chunk = " ".join(words[start:start + size])
            if chunk.strip():
                chunks.append((chunk, page_num))
            start += size - overlap
    return chunks


def chunk_page_text(text: str, size: int = 800, overlap: int = 100) -> list[str]:
    """Chunk a single page's text into a list of chunk strings. Same word
    windowing as chunk_pages but for one page (the page_number is supplied
    by the caller in the page-by-page loop)."""
    words = _sanitize_text(text).split()
    chunks: list[str] = []
    start = 0
    while start < len(words):
        chunk = " ".join(words[start:start + size])
        if chunk.strip():
            chunks.append(chunk)
        start += size - overlap
    return chunks


def classify_content_type(text: str) -> str:
    """Lightweight heuristic: figure if a single [bracketed] note, math if LaTeX markers."""
    stripped = text.strip()
    if stripped.startswith("[") and stripped.endswith("]") and len(stripped) < 500:
        return "figure"
    math_markers = ("\\frac", "\\int", "\\sum", "\\sqrt", "$$", "\\(")
    if any(m in text for m in math_markers) or text.count("$") >= 4:
        return "math"
    return "text"


@transient()
def embed(text: str, ctx=None) -> list[float]:
    res = track_gemini_embed(
        "embed_chunk",
        model="gemini-embedding-001",
        contents=text,
        config={"output_dimensionality": 1536},
        ctx=ctx,
    )
    return res.embeddings[0].values


@transient()
def embed_many(texts: list[str], ctx=None) -> list[list[float]]:
    """Embed a batch of texts in one Gemini call, preserving order.

    Replaces the per-chunk embed() loop in the ingest path: one page's
    chunks all embed together, which cuts request count (and latency) on
    long pages. Keeps output_dimensionality=1536 to match the
    chunks.embedding column, stays @transient for the same 429/transient
    retry behaviour, and returns one vector per input text in the same
    order. An empty input returns an empty list without calling the API."""
    if not texts:
        return []
    res = track_gemini_embed(
        "embed_chunk",
        model="gemini-embedding-001",
        contents=texts,
        config={"output_dimensionality": 1536},
        ctx=ctx,
    )
    return [e.values for e in res.embeddings]


def build_outline(text: str, ctx=None) -> str:
    # Haiku 4.5 handles structured outline extraction well at ~1/3 the
    # cost of Sonnet 4.6. Watch the next few ingests for outline quality
    # (topic granularity, off-topic bleed). Roll back the model string if
    # outlines start missing topics or grouping unrelated sections.
    msg = track_claude(
        "build_outline",
        model="claude-haiku-4-5",
        max_tokens=2000,
        messages=[{"role": "user", "content": OUTLINE_PROMPT + text[:50000] + STYLE_RULES}],
        ctx=ctx,
    )
    return msg.content[0].text


def _set_progress(doc_id: str, text: str):
    """Write a short status string the frontend can display.
    Best-effort: if the DB write fails we log and keep going."""
    try:
        supabase.table("documents").update({"progress": text}) \
            .eq("id", doc_id).execute()
    except Exception as e:
        log.warning("could not write progress for doc_id=%s: %s", doc_id, e)


def _read_ingest_cursor(doc_id: str) -> int:
    """Last fully-completed 1-indexed page from a prior (interrupted)
    ingest, or 0 if none. Resumed ingests skip pages <= this value.
    Best-effort: any read failure starts from 0 (re-doing pages is safe
    because each page's rows are written delete-then-insert)."""
    try:
        rows = supabase.table("documents").select("ingest_cursor") \
            .eq("id", doc_id).execute().data or []
        if rows and rows[0].get("ingest_cursor") is not None:
            return int(rows[0]["ingest_cursor"])
    except Exception as e:
        log.warning("could not read ingest_cursor for doc_id=%s: %s", doc_id, e)
    return 0


def ingest_document(user_id: str, doc_id: str, files: list[tuple[bytes, str]],
                    chapter: str | None = None):
    """Ingest one uploaded document, page by page.

    `files` is a list of (bytes, filename). Normal case is a single file
    (PDF / docx / pptx / txt / md / one image). A multi-image upload
    arrives as several image files that become ONE document whose pages
    are the images in upload order. The dispatch is by the first file's
    extension (the caller has already validated that the set is
    homogeneous per the upload contract).

    `chapter` (optional) is the raw chapter label the user typed (e.g. "5",
    "V", "Chapter 5", "five"). It applies ONLY to a single PDF upload; for
    non-PDF or multi-file uploads it is ignored and the whole document is
    ingested. When given and parseable, only that chapter's page span is
    OCR'd / embedded (cutting API cost) and the chapter number is appended
    to the document title so the library shows the scope. An unparseable
    label, or a chapter that can't be located in the PDF, marks the document
    failed with an actionable message and returns (no partial ingest).

    Idempotency is per page, keyed on (document_id, page_number) via
    delete-then-insert, so an interrupted ingest can resume from
    documents.ingest_cursor without duplicating rows. The global running
    chunk_index is NOT used as the idempotency key — it shifts on
    re-ingest and would not be stable.
    """
    primary_bytes, primary_name = files[0]
    total_bytes = sum(len(b) for b, _ in files)
    try:
        log.info("ingest start doc_id=%s filename=%s files=%d size_bytes=%d",
                 doc_id, primary_name, len(files), total_bytes)
        _set_progress(doc_id, "extracting text")

        ext = _ext(primary_name)
        is_image_set = len(files) > 1 or ext in IMAGE_EXTS

        # Chapter scope applies ONLY to a single PDF. For everything else the
        # request is ignored and the whole document is ingested (unchanged).
        page_range: tuple[int, int] | None = None
        chapter_label = (chapter or "").strip()
        if chapter_label and not is_image_set and ext == ".pdf":
            chapter_n = parse_chapter_request(chapter_label)
            if chapter_n is None:
                log.info("ingest could not parse chapter %r for doc_id=%s",
                         chapter_label, doc_id)
                supabase.table("documents").update({
                    "status": "failed",
                    "progress": CHAPTER_UNPARSEABLE_MESSAGE.format(label=chapter_label),
                }).eq("id", doc_id).execute()
                return
            doc = fitz.open(stream=primary_bytes, filetype="pdf")
            page_range = find_chapter_page_range(doc, chapter_n)
            if page_range is None:
                log.info("ingest could not locate chapter %d for doc_id=%s",
                         chapter_n, doc_id)
                supabase.table("documents").update({
                    "status": "failed",
                    "progress": CHAPTER_NOT_FOUND_MESSAGE.format(n=chapter_n),
                }).eq("id", doc_id).execute()
                return
            log.info("ingest restricting doc_id=%s to chapter %d pages %s",
                     doc_id, chapter_n, page_range)
            # Persist the chapter scope on the row AND annotate the title so
            # the library shows the scope. Persisting `chapter` (the raw label)
            # lets a Retry (POST /reprocess) re-run ingest with the SAME
            # chapter span; without it the retry would fall back to the whole
            # book. Read the current title and append the marker once.
            try:
                rows = supabase.table("documents").select("title") \
                    .eq("id", doc_id).execute().data or []
                cur_title = (rows[0].get("title") if rows else None) or "Document"
                marker = f" — Chapter {chapter_n}"
                patch: dict = {"chapter": chapter_label}
                if not cur_title.endswith(marker):
                    patch["title"] = cur_title + marker
                supabase.table("documents").update(patch) \
                    .eq("id", doc_id).execute()
            except Exception as e:
                log.warning("could not persist chapter scope for doc_id=%s: %s",
                            doc_id, e)

        if is_image_set and ext in IMAGE_EXTS:
            # One or more image files -> ONE document, images as ordered pages.
            pages = extract_image_pages(files, doc_id=doc_id)
            source_type = "image"
            ocr_page_numbers: set[int] = set()
        elif page_range is not None:
            # Chapter scope active: thread the restricted span through.
            pages, source_type, ocr_page_numbers = extract_pages(
                primary_bytes, primary_name, doc_id=doc_id, page_range=page_range)
        else:
            # Whole document (no chapter): keep the call to the narrow
            # (file_bytes, filename, doc_id) form so extract_pages isn't
            # handed a page_range it doesn't need.
            pages, source_type, ocr_page_numbers = extract_pages(
                primary_bytes, primary_name, doc_id=doc_id)
        log.info(
            "ingest extracted %d pages, source_type=%s, ocr_pages=%d for doc_id=%s",
            len(pages), source_type, len(ocr_page_numbers), doc_id)

        if not any(_sanitize_text(text).strip() for _, text in pages):
            log.warning("ingest got no extractable text for doc_id=%s", doc_id)
            supabase.table("documents").update({
                "status": "failed",
                "progress": EMPTY_DOC_MESSAGE,
            }).eq("id", doc_id).execute()
            return

        # Whether the PDF figure-extraction path applies. Multi-image sets
        # and office docs have no embedded-figure stage (each image IS the
        # page; office-doc figures are out of scope — text only).
        figures_from_pdf = (not is_image_set) and ext == ".pdf"

        # Resume support: skip pages already fully written by a prior run.
        cursor = _read_ingest_cursor(doc_id)
        if cursor:
            log.info("ingest resuming doc_id=%s from cursor=%d", doc_id, cursor)

        ctx = {"doc_id": doc_id}
        total_pages = len(pages)
        chunk_index = 0  # running display index only — NOT an idempotency key

        for page_num, page_text in pages:
            if page_num <= cursor:
                # Account for already-written chunks so the running display
                # index stays roughly continuous on resume.
                chunk_index += len(chunk_page_text(page_text))
                continue

            _set_progress(doc_id, f"embedding page {page_num} of {total_pages}")

            page_chunks = chunk_page_text(page_text)

            # Figures for THIS page. Two storage subdirs:
            #   - <user>/<doc>/figures/...  → embedded images from text PDF
            #     pages (PyMuPDF's `page.get_images()` path). Safe.
            #   - <user>/<doc>/diagrams/... → vision-cropped diagrams from
            #     OCR'd PDF pages (`_detect_page_diagrams` path). Safe — only
            #     the diagram region was cropped, not the surrounding text.
            # The runtime "is this safe to show on a test?" check
            # distinguishes the two by path prefix.
            figure_paths: list[str] = []
            if figures_from_pdf:
                is_ocr_page = page_num in ocr_page_numbers
                subdir = "diagrams" if is_ocr_page else "figures"
                try:
                    pngs = _page_figures(
                        primary_bytes, primary_name, page_num,
                        is_ocr_page, page_text, doc_id=doc_id)
                    for idx, png in enumerate(pngs):
                        fp = f"{user_id}/{doc_id}/{subdir}/p{page_num}_{idx}.png"
                        supabase.storage.from_("uploads").upload(
                            fp, png,
                            {"content-type": "image/png", "upsert": "true"})
                        figure_paths.append(fp)
                    if figure_paths:
                        log.info("ingest uploaded %d figure images for page=%s "
                                 "doc_id=%s", len(figure_paths), page_num, doc_id)
                except Exception:
                    log.exception("figure image upload failed for page=%s "
                                  "doc_id=%s", page_num, doc_id)

            # Batched embedding: every chunk on this page in ONE call.
            embeddings = embed_many(page_chunks, ctx=ctx)

            # Build this page's chunk rows. Each text chunk claims the next
            # figure on the page in order (figure_cursor); leftover figures
            # become orphan figure-only rows in THIS page's batch.
            rows = []
            figure_cursor = 0
            for chunk, emb in zip(page_chunks, embeddings):
                ctype = classify_content_type(chunk)
                figure_path = None
                if figure_cursor < len(figure_paths):
                    figure_path = figure_paths[figure_cursor]
                    figure_cursor += 1
                rows.append({
                    "document_id": doc_id,
                    "user_id": user_id,
                    "content": chunk,
                    "embedding": emb,
                    "chunk_index": chunk_index,
                    "page_number": page_num,
                    "content_type": ctype,
                    "figure_path": figure_path,
                })
                chunk_index += 1

            # Backfill orphan figures for THIS page. A page can have more
            # extracted images than text chunks (e.g. a composite figure
            # with 4 subfigures but one caption chunk). Each leftover gets a
            # dedicated figure-only row (empty content + placeholder
            # embedding) so all subfigures surface on the lesson screen
            # while staying out of vector search. These belong to THIS
            # page's batch, so re-ingesting the page rewrites them too.
            if figure_cursor < len(figure_paths):
                placeholder_emb = embed("figure", ctx=ctx)
                for fp in figure_paths[figure_cursor:]:
                    rows.append({
                        "document_id": doc_id,
                        "user_id": user_id,
                        "content": "",
                        "embedding": placeholder_emb,
                        "chunk_index": chunk_index,
                        "page_number": page_num,
                        "content_type": "figure",
                        "figure_path": fp,
                    })
                    chunk_index += 1

            # Per-page delete-then-insert keyed on (document_id, page_number).
            # This is the idempotency boundary: re-running a page wipes its
            # old rows (including its orphan figure rows) before inserting
            # the fresh batch, so a resumed ingest never duplicates.
            supabase.table("chunks").delete() \
                .eq("document_id", doc_id).eq("page_number", page_num).execute()
            if rows:
                supabase.table("chunks").insert(rows).execute()
            log.info("ingest wrote %d chunk rows for page=%s doc_id=%s",
                     len(rows), page_num, doc_id)

            # Mark the page fully complete so a crash after this resumes
            # from the NEXT page.
            supabase.table("documents").update({"ingest_cursor": page_num}) \
                .eq("id", doc_id).execute()

        # Outline is built from the full text AFTER all pages are written.
        full_text = _sanitize_text("\n\n".join(text for _, text in pages))
        log.info("ingest building outline for doc_id=%s (%d chars)",
                 doc_id, len(full_text))
        _set_progress(doc_id, "building outline")
        outline = _sanitize_text(build_outline(full_text, ctx=ctx))

        supabase.table("documents").update({
            "source_type": source_type,
            "outline": outline,
            "status": "ready",
            "progress": None,
        }).eq("id", doc_id).execute()
        log.info("ingest complete doc_id=%s", doc_id)

    except QuotaExhausted as e:
        log.warning("ingestion hit Gemini quota for doc_id=%s: %s", doc_id, e)
        supabase.table("documents").update({
            "status": "failed",
            "progress": QUOTA_MESSAGE,
        }).eq("id", doc_id).execute()
        # Don't re-raise: this is a known operational state, not a code bug.
        # ingest_cursor is left in place so a retry resumes mid-document.
    except Exception:
        log.exception("ingestion failed for doc_id=%s user_id=%s", doc_id, user_id)
        # Leave the last `progress` string in place so the failure point is visible
        supabase.table("documents").update({"status": "failed"}).eq("id", doc_id).execute()
        raise
