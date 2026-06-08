"""Format-specific text extractors:
  .txt / .md  -> utf-8 decode, single page
  .docx       -> paragraphs joined, single page (built in-memory)
  .pptx       -> one slide per page (built in-memory)
  images      -> ordered pages via OCR (OCR monkeypatched, deterministic)
And the extract_pages() extension dispatch.
"""
import io

from app import ingest


# --------------------------------------------------------------------------
# .txt / .md
# --------------------------------------------------------------------------

def test_extract_text_file_decodes_utf8():
    pages = ingest.extract_text_file("hello\nworld ✓".encode("utf-8"))
    assert pages == [(1, "hello\nworld ✓")]


def test_extract_text_file_replaces_bad_bytes():
    # invalid utf-8 byte -> replacement char, no crash, still one page
    pages = ingest.extract_text_file(b"ok\xff end")
    assert len(pages) == 1
    assert pages[0][0] == 1
    assert "ok" in pages[0][1] and "end" in pages[0][1]


def test_extract_pages_dispatches_txt_and_md():
    p_txt, st_txt, ocr_txt = ingest.extract_pages(b"plain", "notes.txt")
    p_md, st_md, ocr_md = ingest.extract_pages(b"# head", "README.md")
    assert p_txt == [(1, "plain")] and st_txt == "text" and ocr_txt == set()
    assert p_md == [(1, "# head")] and st_md == "text" and ocr_md == set()


# --------------------------------------------------------------------------
# .docx — built in-memory with python-docx
# --------------------------------------------------------------------------

def _make_docx(paragraphs):
    import docx
    d = docx.Document()
    for p in paragraphs:
        d.add_paragraph(p)
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


def test_extract_docx_joins_nonempty_paragraphs():
    data = _make_docx(["First paragraph.", "", "   ", "Second paragraph."])
    pages = ingest.extract_docx(data)
    assert len(pages) == 1
    assert pages[0][0] == 1
    text = pages[0][1]
    assert "First paragraph." in text
    assert "Second paragraph." in text
    # blank / whitespace-only paragraphs are dropped
    assert text == "First paragraph.\nSecond paragraph."


def test_extract_pages_dispatches_docx():
    data = _make_docx(["Body text here."])
    pages, source_type, ocr = ingest.extract_pages(data, "essay.docx")
    assert source_type == "docx"
    assert ocr == set()
    assert "Body text here." in pages[0][1]


# --------------------------------------------------------------------------
# .pptx — one slide per page, built in-memory with python-pptx
# --------------------------------------------------------------------------

def _make_pptx(slides_text):
    """slides_text: list of (title, body) per slide."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    blank = prs.slide_layouts[6]  # blank layout, we add our own textboxes
    for title, body in slides_text:
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
        tf = tb.text_frame
        tf.text = title
        if body:
            tf.add_paragraph().text = body
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_extract_pptx_one_page_per_slide():
    data = _make_pptx([("Slide One Title", "slide one body"),
                       ("Slide Two Title", "slide two body")])
    pages = ingest.extract_pptx(data)
    assert [pn for pn, _ in pages] == [1, 2]
    assert "Slide One Title" in pages[0][1]
    assert "slide one body" in pages[0][1]
    assert "Slide Two Title" in pages[1][1]
    assert "slide two body" in pages[1][1]
    # page 1's text must not leak into page 2
    assert "Slide One Title" not in pages[1][1]


def test_extract_pages_dispatches_pptx():
    data = _make_pptx([("Only Slide", "content")])
    pages, source_type, ocr = ingest.extract_pages(data, "deck.pptx")
    assert source_type == "pptx"
    assert ocr == set()
    assert len(pages) == 1


# --------------------------------------------------------------------------
# Multiple images -> ordered pages (OCR monkeypatched, deterministic)
# --------------------------------------------------------------------------

def test_extract_image_pages_ordered(monkeypatch):
    # OCR returns text derived from the image bytes so order is observable.
    def fake_read_image(img_bytes, model=ingest.MODEL_FAST, ctx=None):
        return "OCR:" + img_bytes.decode()

    monkeypatch.setattr(ingest, "read_image", fake_read_image)
    files = [(b"img-a", "a.png"), (b"img-b", "b.png"), (b"img-c", "c.png")]
    pages = ingest.extract_image_pages(files, doc_id="d1")
    assert pages == [(1, "OCR:img-a"), (2, "OCR:img-b"), (3, "OCR:img-c")]


def test_extract_pages_single_image_dispatch(monkeypatch):
    calls = {}

    def fake_read_image(img_bytes, model=ingest.MODEL_FAST, ctx=None):
        calls["ctx"] = ctx
        return "single-image-text"

    monkeypatch.setattr(ingest, "read_image", fake_read_image)
    pages, source_type, ocr = ingest.extract_pages(b"img", "scan.png", doc_id="d9")
    assert pages == [(1, "single-image-text")]
    assert source_type == "image"
    assert ocr == set()
    # the single-image path passes doc_id through as ctx
    assert calls["ctx"] == {"doc_id": "d9"}


# --------------------------------------------------------------------------
# _ext helper
# --------------------------------------------------------------------------

def test_ext_helper():
    assert ingest._ext("a.PDF") == ".pdf"
    assert ingest._ext("photo.JPG") == ".jpg"
    assert ingest._ext("noext") == ""
    assert ingest._ext("archive.tar.gz") == ".gz"
